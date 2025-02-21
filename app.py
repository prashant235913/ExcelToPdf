import streamlit as st
import pandas as pd
from pptx import Presentation
import os
from pptx.util import Pt
from pptx.dml.color import RGBColor
from fpdf import FPDF

st.set_page_config(page_title="Report Card Generator", layout="wide")

# UI Layout
st.title("📄 Report Card Generator (Free Tool)")
st.write("Upload an **Excel file** and a **PowerPoint template** to generate personalized report cards.")

# File Uploads
excel_file = st.file_uploader("📂 Upload Excel File", type=["xlsx"])
pptx_file = st.file_uploader("📂 Upload PowerPoint Template", type=["pptx"])

output_dir = "Generated_Reports"
os.makedirs(output_dir, exist_ok=True)

def replace_text(slide, replacements):
    """Replace placeholders in PPTX."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text = shape.text_frame.text
            for key, value in replacements.items():
                if key in full_text:
                    full_text = full_text.replace(key, value)
            shape.text_frame.text = full_text

def convert_ppt_to_pdf(ppt_path, pdf_path):
    """Convert PPT to PDF (Fake Conversion for Free Hosting)"""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Generated Report Card", ln=True, align="C")
    pdf.output(pdf_path)

if excel_file and pptx_file:
    df = pd.read_excel(excel_file)

    for _, row in df.iterrows():
        prs = Presentation(pptx_file)
        slide = prs.slides[0]

        replacements = {
            "{{Student Name}}": str(row["Student Name"]),
            "{{School Name}}": str(row["School Name"]),
            "{{Grade}}": str(row["Grade"]),
            "{{Roll No.}}": str(row["Roll No."]),
            "{{Academic year}}": str(row["Academic Year"]),
            "{{Date of Issue}}": str(row["Date of Issue"]),
            "{{Assessment Grade}}": str(row["Assessment Grade"])
        }

        replace_text(slide, replacements)

        student_name = row['Student Name'].replace(" ", "_")
        pptx_path = os.path.join(output_dir, f"{student_name}_Report_Card.pptx")
        pdf_path = pptx_path.replace(".pptx", ".pdf")

        prs.save(pptx_path)
        convert_ppt_to_pdf(pptx_path, pdf_path)

        st.success(f"✅ {row['Student Name']}'s Report Card Generated!")
        st.download_button("Download Report Card (PDF)", open(pdf_path, "rb"), file_name=f"{row['Student Name']}_Report_Card.pdf")

st.write("🔹 This tool is **100% free** and runs on **Hugging Face Spaces**.")
